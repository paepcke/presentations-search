#!/usr/bin/env python
# -*- coding: utf-8 -*-
# @Author: Andreas Paepcke
# @Date:   2026-03-28 15:58:06
# @Last Modified by:   Andreas Paepcke
# @Last Modified time: 2026-03-29 19:14:35

"""
presentation_indexer.py  --  Index Keynote (.key) and PowerPoint (.pptx) files
into a Qdrant vector store for semantic search.

Each presentation is chunked into per-slide documents carrying:
  * slide title
  * body text (all text boxes / bullet runs)
  * speaker notes
  * file path, slide index, and content type as metadata

Keynote files are unpacked via ``keynote-parser`` into a temporary directory;
text is extracted by walking the resulting YAML archives.  PowerPoint files are
handled directly by ``python-pptx``.

Only changed or new files are re-indexed on subsequent runs (mtime manifest).
Deleted files are pruned from the collection automatically.

Usage
-----
    python presentation_indexer.py presentations/ ~/Dropbox/Talks MySlides.key
    python presentation_indexer.py ~/Talks --no-recursive
    python presentation_indexer.py ~/Talks --force --verbose

Options
-------
    inputs             One or more directories or .key/.pptx files.
    --index-dir  DIR   Qdrant storage + manifest directory.
                       Default: ~/.presentation_index
    --collection NAME  Qdrant collection name.  Default: presentation_index
    --batch-size N     Embedding batch size.  Default: 32
    --force            Ignore mtime manifest; re-index everything.
    --no-recursive     Only look one level deep in directories (no descent).
    --verbose          Log each file path as it is processed.
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import io
import json
import re
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Generator

import yaml
from logging_service import LoggingService
import requests
from qdrant_client import QdrantClient
from qdrant_client.models import (
    Distance,
    FieldCondition,
    Filter,
    MatchText,
    MatchValue,
    PointStruct,
    VectorParams,
)

try:
    from PIL import Image as PILImage
    _PIL_AVAILABLE = True
except ImportError:
    _PIL_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    from pptx.enum.text import PP_ALIGN  # noqa: F401
    from pptx.util import Pt             # noqa: F401
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

# ---------------------------------------------------------------------------
# Module-level logger (mirrors code-search project convention)
# ---------------------------------------------------------------------------
log = LoggingService()

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
DEFAULT_INDEX_DIR  = Path("~/.presentation_index")
COLLECTION_NAME    = "presentation_index"
OLLAMA_EMBED_MODEL = "nomic-embed-text"
EMBEDDING_DIM      = 768
DEFAULT_OLLAMA_URL = "http://localhost:11434"
SUPPORTED_SUFFIXES = {".key", ".pptx"}

# Protobuf type names that carry actual slide text in Keynote YAML
_STORAGE_PBTYPE = "TSWP.StorageArchive"
# Object-replacement character used as image/table placeholder in text runs
_ORC = "\uFFFC"


# ---------------------------------------------------------------------------
# SlideChunk dataclass (plain dict so it serialises easily into Qdrant payload)
# ---------------------------------------------------------------------------

def _make_chunk(
    file_path: str,
    slide_index: int,
    slide_position: int,
    title: str,
    body: str,
    notes: str,
    source_format: str,
    image_only: bool = False,
) -> dict:
    """Return a payload dict for a single slide chunk.

    :param file_path:      Absolute path to the source presentation file.
    :param slide_index:    0-based index derived from sorted YAML filename order.
                           Not meaningful as a human-readable slide number.
    :param slide_position: 1-based slide position matching the presentation
                           application's numbering (derived from slideTree order
                           for Keynote; from slide order for PPTX).
    :param title:          Slide title text (may be empty).
    :param body:           All non-title, non-note text joined with newlines.
    :param notes:          Speaker-note text (may be empty).
    :param source_format:  ``"keynote"`` or ``"pptx"``.
    :param image_only:     ``True`` if no text was extractable (image/diagram slide).
    :return: Payload dict suitable for a Qdrant ``PointStruct``.
    """
    return {
        "file_path":      file_path,
        "slide_index":    slide_index,
        "slide_position": slide_position,
        "title":          title,
        "body":           body,
        "notes":          notes,
        "source_format":  source_format,
        "image_only":     image_only,
        # Combined text field used for embedding
        "text": "\n".join(filter(None, [title, body, notes])),
    }


OLLAMA_VISION_MODEL = "llama3.2-vision"

# Prompt sent to the vision model for each image-only slide.
# Asks for visible text first (most useful for search), then content description.
_VISION_PROMPT = (
    "This is a presentation slide. "
    "First, transcribe any text you can see verbatim. "
    "Then describe what is depicted: charts, diagrams, photos, logos, or other visual content. "
    "Include any axis labels, legend entries, or callout text. "
    "Be concise and factual."
)


# ---------------------------------------------------------------------------
# SlideVisionCaptioner
# ---------------------------------------------------------------------------

class SlideVisionCaptioner:
    """Generate text captions for image-only slides via Ollama's vision model.

    Sends the slide image to ``llama3.2-vision`` via the Ollama chat endpoint
    and returns a text description suitable for embedding and search.

    :param ollama_url: Base URL for the Ollama REST API.
    """

    def __init__(self, ollama_url: str = DEFAULT_OLLAMA_URL) -> None:
        self._chat_url = f"{ollama_url.rstrip('/')}/api/chat"

    # ------------------------------------------------------------------
    def caption(self, image_path: Path) -> str:
        """Generate a caption for the image at *image_path*.

        :param image_path: Path to a PNG or JPEG image file.
        :return: Caption string, or empty string on failure.
        :raises RuntimeError: If the Ollama API call fails.
        """
        try:
            raw_bytes = image_path.read_bytes()
        except Exception as exc:
            log.warn(f"Could not read image '{image_path}': {exc}")
            return ""

        # Resize to max 1024px on the longest side before encoding — large
        # high-resolution images cause vision model timeouts with no quality
        # benefit for text/label extraction.
        if _PIL_AVAILABLE:
            try:
                img     = PILImage.open(io.BytesIO(raw_bytes))
                max_dim = max(img.width, img.height)
                if max_dim > 1024:
                    scale    = 1024 / max_dim
                    new_size = (int(img.width * scale), int(img.height * scale))
                    img      = img.resize(new_size, PILImage.LANCZOS)
                buf = io.BytesIO()
                img.save(buf, format=img.format or "PNG")
                raw_bytes = buf.getvalue()
            except Exception as exc:
                log.warn(f"Could not resize '{image_path.name}': {exc} — using original")
        else:
            log.warn("Pillow not installed — sending full-size image to vision model")

        b64 = base64.b64encode(raw_bytes).decode()

        payload = {
            "model":  OLLAMA_VISION_MODEL,
            "stream": False,
            "messages": [{
                "role":    "user",
                "content": _VISION_PROMPT,
                "images":  [b64],
            }],
        }

        for attempt in range(2):   # one retry on timeout
            try:
                resp = requests.post(
                    self._chat_url, json=payload, timeout=300
                )
                resp.raise_for_status()
                return resp.json()["message"]["content"].strip()
            except requests.exceptions.Timeout:
                if attempt == 0:
                    log.warn(f"Vision timeout on '{image_path.name}' — retrying …")
                    continue
                raise RuntimeError(
                    f"Ollama vision API timed out twice for '{image_path}'"
                )
            except Exception as exc:
                raise RuntimeError(
                    f"Ollama vision API failed for '{image_path}': {exc}"
                ) from exc


# ---------------------------------------------------------------------------
# Keynote extractor
# ---------------------------------------------------------------------------

class KeynoteExtractor:
    """Extract per-slide text chunks from a Keynote ``.key`` file.

    Uses ``keynote-parser`` to unpack the file into a temporary directory of
    YAML archives, then walks each ``Slide-*.iwa.yaml`` file to collect
    ``TSWP.StorageArchive`` objects, distinguishing title placeholders,
    body text, and speaker notes.

    :param key_path: Path to the ``.key`` file to extract.
    """

    def __init__(
        self,
        key_path:  Path,
        captioner: "SlideVisionCaptioner | None" = None,
        verbose:   bool = False,
    ) -> None:
        self._key_path  = key_path
        self._captioner = captioner
        self._verbose   = verbose

    # ------------------------------------------------------------------
    def extract(self) -> list[dict]:
        """Unpack the Keynote file and return a list of slide chunk dicts.

        Creates a temporary directory for unpacking that is cleaned up
        automatically after extraction.  If a ``captioner`` was supplied,
        vision captioning runs *inside* the temp dir context so that image
        files are still accessible.

        :return: List of slide payload dicts (see :func:`_make_chunk`).
        :raises RuntimeError: If ``keynote-parser`` is not installed or
            the unpack command fails.
        """
        chunks: list[dict] = []

        with tempfile.TemporaryDirectory(prefix="kn_unpack_") as tmp_dir:
            tmp_path   = Path(tmp_dir)
            unpack_dir = tmp_path / "unpacked"

            result = subprocess.run(
                [
                    "keynote-parser", "unpack",
                    str(self._key_path),
                    "--output", str(unpack_dir),
                ],
                capture_output=True,
                text=True,
            )
            if result.returncode != 0:
                raise RuntimeError(
                    f"keynote-parser failed on '{self._key_path}': "
                    f"{result.stderr.strip()}"
                )

            index_dir = unpack_dir / "Index"
            if not index_dir.is_dir():
                log.warn(
                    f"No Index/ directory after unpacking '{self._key_path}' "
                    f"— skipping."
                )
                return chunks

            data_dir = unpack_dir / "Data"

            # Build archive-id → 1-based Keynote position map
            position_map, total_positions = self._build_position_map(index_dir)

            # Build Data/ id → filename lookup once for all slides
            asset_map = self._build_asset_map(data_dir)

            # image_map: slide_position → Path (only populated when captioning)
            image_map: dict[int, Path] = {}

            # Collect slide YAML files, sort by name for stable slide order.
            # Modern Keynote produces individual Slide-XXXXXX.iwa.yaml files;
            # older formats store all slides in a single Slide.iwa.yaml.
            slide_files = sorted(index_dir.glob("Slide-*.iwa.yaml"))
            if not slide_files:
                singular = index_dir / "Slide.iwa.yaml"
                if singular.exists():
                    slide_files = [singular]
            for slide_idx, yaml_file in enumerate(slide_files):
                # Extract archive id from filename: Slide-XXXXXX.iwa.yaml
                # Falls back to -1 for the singular Slide.iwa.yaml format.
                stem = yaml_file.stem  # e.g. "Slide-106936.iwa" or "Slide.iwa"
                try:
                    archive_id = int(stem.split("Slide-")[1].replace(".iwa", ""))
                except (IndexError, ValueError):
                    archive_id = -1
                slide_pos    = position_map.get(archive_id, slide_idx + 1)
                slide_chunks = self._parse_slide_yaml(yaml_file, slide_idx, slide_pos)
                chunks.extend(slide_chunks)

                # For image-only slides, find the best asset image
                if self._captioner and slide_chunks and slide_chunks[0].get("image_only"):
                    img = self._find_slide_image(yaml_file, asset_map)
                    if img:
                        image_map[slide_pos] = img

            # Emit image-only stubs for slideTree positions with no YAML file.
            # These slides (e.g. images imported from Affinity Photo) cannot be
            # captioned reliably — the only available fallback image is the deck
            # cover (preview.jpg), which shows a different slide entirely.
            # Leave them as stubs; Option C (vision captioning) only applies to
            # slides whose asset can be identified from their YAML.
            mapped_positions = set(position_map.values())

            for pos in range(1, total_positions + 1):
                if pos not in mapped_positions:
                    chunks.append(_make_chunk(
                        file_path=str(self._key_path.resolve()),
                        slide_index=-1,
                        slide_position=pos,
                        title="[image-only slide]",
                        body="",
                        notes="",
                        source_format="keynote",
                        image_only=True,
                    ))

            # Apply captioning NOW while temp dir (and images) still exist
            if self._captioner and image_map:
                chunks = self._apply_captions(chunks, image_map)

        return chunks

    # ------------------------------------------------------------------
    def _apply_captions(
        self,
        chunks:    list[dict],
        image_map: dict[int, Path],
    ) -> list[dict]:
        """Replace image-only stub titles with vision model captions.

        Called while the temp directory is still live so image paths are valid.

        :param chunks:    Slide chunk dicts from extraction.
        :param image_map: Mapping of ``slide_position`` → image ``Path``.
        :return: Updated chunk list.
        """
        total = sum(1 for c in chunks if c.get("image_only") and
                    c.get("slide_position") in image_map)
        done  = 0
        log.info(f"Captioning {total} image-only slide(s) via vision model …")

        for chunk in chunks:
            if not chunk.get("image_only"):
                continue
            pos = chunk.get("slide_position")
            img = image_map.get(pos)
            if not img:
                continue
            try:
                caption = self._captioner.caption(img)
                if caption:
                    chunk["title"]      = caption[:200]
                    chunk["body"]       = caption
                    chunk["text"]       = caption
                    chunk["image_only"] = False
                    done += 1
                    if self._verbose:
                        log.info(f"  Slide {pos}: captioned ({len(caption)} chars)")
            except Exception as exc:
                log.warn(f"  Slide {pos}: captioning failed — {exc}")

        log.info(f"Captioned {done}/{total} image-only slide(s).")
        return chunks

    # ------------------------------------------------------------------
    @staticmethod
    def _build_asset_map(data_dir: Path) -> dict[str, Path]:
        """Build a mapping from asset identifier to file path in ``Data/``.

        Keynote asset filenames end with ``-<identifier>.<ext>`` where
        identifier matches the ``TSD.ImageArchive.data.identifier`` field.
        Prefers full-size files over ``-small-`` variants.

        :param data_dir: Path to the unpacked ``Data/`` directory.
        :return: Dict mapping identifier string → ``Path``.
        """
        if not data_dir.is_dir():
            return {}

        asset_map: dict[str, Path] = {}
        for fpath in data_dir.iterdir():
            m = re.search(r'-(\d+)\.(\w+)$', fpath.name)
            if not m:
                continue
            asset_id = m.group(1)
            # Prefer full-size over small variants — only overwrite if not set
            # or current entry is a small variant
            existing = asset_map.get(asset_id)
            if existing is None or "-small-" in existing.name:
                asset_map[asset_id] = fpath
        return asset_map

    # ------------------------------------------------------------------
    @staticmethod
    def _find_slide_image(yaml_file: Path, asset_map: dict[str, Path]) -> Path | None:
        """Find the best image asset for an image-only slide.

        Walks the slide YAML for a ``TSD.ImageArchive`` object and extracts
        its ``data.identifier``, then looks it up in *asset_map*.

        :param yaml_file:  Path to the slide's unpacked YAML file.
        :param asset_map:  Mapping from asset id string to ``Path``.
        :return: ``Path`` to the image file, or ``None`` if not found.
        """
        try:
            with open(yaml_file, encoding="utf-8") as fh:
                data = yaml.safe_load(fh)
        except Exception:
            return None

        for chunk in data.get("chunks", []):
            for archive in chunk.get("archives", []):
                for obj in archive.get("objects", []):
                    if obj.get("_pbtype") == "TSD.ImageArchive":
                        data_ref = obj.get("data", {})
                        if isinstance(data_ref, dict):
                            asset_id = str(data_ref.get("identifier", ""))
                            if asset_id and asset_id in asset_map:
                                return asset_map[asset_id]
        return None

    # ------------------------------------------------------------------
    @staticmethod
    def _build_position_map(index_dir: Path) -> dict[int, int]:
        """Parse ``Document.iwa.yaml`` to map YAML archive ids to 1-based
        Keynote slide positions.

        The ``KN.ShowArchive`` in ``Document.iwa.yaml`` contains a
        ``slideTree.slides`` list giving the presentation order.  Each entry
        references a slideTree identifier that is slightly larger than the
        corresponding ``Slide-XXXXXX.iwa.yaml`` archive id.  We match each
        slideTree id to the largest yaml archive id that is ≤ the tree id
        with a delta ≤ ``_MAX_ID_DELTA`` (empirically ~40; we use 100 for
        safety).

        :param index_dir: Path to the unpacked ``Index/`` directory.
        :return: Dict mapping yaml archive id (int) → 1-based slide position.
        """
        _MAX_ID_DELTA = 100

        doc_yaml = index_dir / "Document.iwa.yaml"
        if not doc_yaml.exists():
            return {}, 0

        try:
            with open(doc_yaml, encoding="utf-8") as fh:
                data = yaml.safe_load(fh)
        except Exception as exc:
            log.warn(f"Could not parse Document.iwa.yaml: {exc} — slide positions unavailable.")
            return {}, 0

        # Locate KN.ShowArchive (referenced as 'show' from KN.DocumentArchive)
        # Find its identifier first
        show_id: str | None = None
        for chunk in data.get("chunks", []):
            for archive in chunk.get("archives", []):
                for obj in archive.get("objects", []):
                    if obj.get("_pbtype") == "KN.DocumentArchive":
                        show_id = obj.get("show", {}).get("identifier")
                        break

        if not show_id:
            return {}, 0

        # Now find the KN.ShowArchive with that identifier
        ordered_tree_ids: list[int] = []
        for chunk in data.get("chunks", []):
            for archive in chunk.get("archives", []):
                if archive.get("header", {}).get("identifier") == show_id:
                    for obj in archive.get("objects", []):
                        slides = obj.get("slideTree", {}).get("slides", [])
                        ordered_tree_ids = [int(s["identifier"]) for s in slides]
                        break

        if not ordered_tree_ids:
            return {}, 0

        # Collect all yaml archive ids from Slide-*.iwa.yaml filenames
        yaml_ids = sorted(
            int(f.stem.split("Slide-")[1].replace(".iwa", ""))
            for f in index_dir.glob("Slide-*.iwa.yaml")
        )

        # Match each tree id to the largest yaml id ≤ tree id with delta ≤ MAX.
        # Each yaml_id may only be claimed once (first match in slideTree order
        # wins).  Tree ids with no unclaimed candidate have no YAML file —
        # keynote-parser dropped them — and are silently skipped.
        position_map: dict[int, int] = {}
        claimed: set[int] = set()

        for position, tree_id in enumerate(ordered_tree_ids, 1):
            candidates = [y for y in yaml_ids
                          if y <= tree_id and (tree_id - y) <= _MAX_ID_DELTA
                          and y not in claimed]
            if candidates:
                matched_yaml_id = max(candidates)
                position_map[matched_yaml_id] = position
                claimed.add(matched_yaml_id)

        return position_map, len(ordered_tree_ids)

    # ------------------------------------------------------------------
    def _parse_slide_yaml(self, yaml_file: Path, slide_idx: int, slide_position: int) -> list[dict]:
        """Parse one ``Slide-XXX.iwa.yaml`` file into zero or more chunks.

        A single YAML file contains one logical slide.  We collect one
        combined chunk per slide with separate title, body, and notes fields.

        :param yaml_file:      Path to the unpacked YAML file.
        :param slide_idx:      0-based index from sorted filename order.
        :param slide_position: 1-based Keynote slide number from slideTree.
        :return: List with 0 or 1 chunk dicts for this slide.
        """
        try:
            with open(yaml_file, encoding="utf-8") as fh:
                data = yaml.safe_load(fh)
        except Exception as exc:
            log.warn(f"Could not parse '{yaml_file}': {exc} — skipping slide.")
            return []

        title_texts: list[str] = []
        body_texts:  list[str] = []
        note_texts:  list[str] = []

        # Walk all archive objects in the chunk list
        for chunk in data.get("chunks", []):
            for archive in chunk.get("archives", []):
                for obj in archive.get("objects", []):
                    if obj.get("_pbtype") != _STORAGE_PBTYPE:
                        continue

                    raw_texts = [
                        t for t in obj.get("text", [])
                        if isinstance(t, str) and t.strip(_ORC).strip()
                    ]
                    if not raw_texts:
                        continue

                    combined = "\n".join(raw_texts).replace(_ORC, "").strip()
                    if not combined:
                        continue

                    kind = obj.get("kind", "")
                    if kind == "NOTE":
                        note_texts.append(combined)
                    else:
                        # Treat the first non-note, single-line short text as
                        # the title.  Multi-line text (newlines present) is
                        # body content even if it appears first.
                        single_line = combined.replace("\n", " ").strip()
                        if not title_texts and "\n" not in combined and len(single_line) < 200:
                            title_texts.append(single_line)
                        else:
                            body_texts.append(combined)

        title = " ".join(title_texts)
        body  = "\n".join(body_texts)
        notes = "\n".join(note_texts)

        if not any([title, body, notes]):
            # Slide contains no extractable text (image/diagram only).
            # Index a stub so the slide is counted and findable by position.
            return [_make_chunk(
                file_path=str(self._key_path.resolve()),
                slide_index=slide_idx,
                slide_position=slide_position,
                title="[image-only slide]",
                body="",
                notes="",
                source_format="keynote",
                image_only=True,
            )]

        return [_make_chunk(
            file_path=str(self._key_path.resolve()),
            slide_index=slide_idx,
            slide_position=slide_position,
            title=title,
            body=body,
            notes=notes,
            source_format="keynote",
        )]


# ---------------------------------------------------------------------------
# PowerPoint extractor
# ---------------------------------------------------------------------------

class PptxExtractor:
    """Extract per-slide text chunks from a PowerPoint ``.pptx`` file.

    Uses ``python-pptx`` to iterate over slides, collecting title placeholder
    text, all other shape text, and speaker notes.

    :param pptx_path: Path to the ``.pptx`` file to extract.
    """

    def __init__(self, pptx_path: Path) -> None:
        self._pptx_path = pptx_path

    # ------------------------------------------------------------------
    def extract(self) -> list[dict]:
        """Return a list of per-slide chunk dicts for the presentation.

        :return: List of slide payload dicts (see :func:`_make_chunk`).
        :raises RuntimeError: If ``python-pptx`` is not installed or the
            file cannot be opened.
        """
        if not _PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx is not installed. Run: pip install python-pptx"
            )
        try:
            prs = PptxPresentation(str(self._pptx_path))
        except Exception as exc:
            raise RuntimeError(
                f"python-pptx could not open '{self._pptx_path}': {exc}"
            ) from exc

        chunks: list[dict] = []

        for slide_idx, slide in enumerate(prs.slides):
            title = ""
            body_parts: list[str] = []
            notes = ""

            # ---- slide body shapes ----
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # TITLE = 1, CENTER_TITLE = 3
                    if ph_type in (1, 3) and not title:
                        title = text
                        continue
                body_parts.append(text)

            # ---- speaker notes ----
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes = notes_tf.text.strip() if notes_tf else ""

            body = "\n".join(body_parts)

            if not any([title, body, notes]):
                # Slide contains no extractable text (image/diagram only).
                chunks.append(_make_chunk(
                    file_path=str(self._pptx_path.resolve()),
                    slide_index=slide_idx,
                    slide_position=slide_idx + 1,
                    title="[image-only slide]",
                    body="",
                    notes="",
                    source_format="pptx",
                    image_only=True,
                ))
                continue

            chunks.append(_make_chunk(
                file_path=str(self._pptx_path.resolve()),
                slide_index=slide_idx,
                slide_position=slide_idx + 1,
                title=title,
                body=body,
                notes=notes,
                source_format="pptx",
            ))

        return chunks


# ---------------------------------------------------------------------------
# Mtime manifest  (identical pattern to code-search project)
# ---------------------------------------------------------------------------

class MtimeManifest:
    """Persist file modification times to detect changed presentations.

    Stores a JSON mapping of ``absolute_path -> mtime_float`` next to the
    Qdrant storage directory.

    :param manifest_path: Path to the JSON manifest file.
    """

    def __init__(self, manifest_path: Path) -> None:
        self._path = manifest_path
        self._data: dict[str, float] = {}
        if manifest_path.exists():
            try:
                self._data = json.loads(manifest_path.read_text())
            except Exception as exc:
                log.warn(f"Could not read manifest '{manifest_path}': {exc} — starting fresh.")

    # ------------------------------------------------------------------
    def is_stale(self, path: Path) -> bool:
        """Return ``True`` if *path* is new or has been modified since last index.

        :param path: Absolute path to the file to check.
        :return: ``True`` if the file should be (re-)indexed.
        """
        key   = str(path.resolve())
        mtime = path.stat().st_mtime
        return self._data.get(key) != mtime

    # ------------------------------------------------------------------
    def mark(self, path: Path) -> None:
        """Record the current mtime for *path*.

        :param path: Absolute path whose mtime should be saved.
        """
        self._data[str(path.resolve())] = path.stat().st_mtime

    # ------------------------------------------------------------------
    def remove(self, path: Path) -> None:
        """Remove *path* from the manifest (called when a file is deleted).

        :param path: Absolute path to remove.
        """
        self._data.pop(str(path.resolve()), None)

    # ------------------------------------------------------------------
    def known_paths(self) -> set[str]:
        """Return all absolute path strings currently tracked in the manifest.

        :return: Set of path strings.
        """
        return set(self._data.keys())

    # ------------------------------------------------------------------
    def save(self) -> None:
        """Flush the manifest to disk.

        Creates parent directories if they do not exist.
        """
        self._path.parent.mkdir(parents=True, exist_ok=True)
        self._path.write_text(json.dumps(self._data, indent=2))


# ---------------------------------------------------------------------------
# Main indexer
# ---------------------------------------------------------------------------

class PresentationIndexer:
    """Orchestrate discovery, extraction, embedding, and Qdrant upsert.

    Walks the supplied input paths (files or directories), extracts slide
    chunks from each supported presentation file, embeds the text with
    ``nomic-embed-text``, and upserts into a local Qdrant collection.

    :param index_dir:      Directory for Qdrant storage and the mtime manifest.
    :param collection:     Qdrant collection name.
    :param batch_size:     Number of chunks to embed and upsert per batch.
    :param force:          If ``True``, ignore the mtime manifest and re-index all.
    :param recursive:      If ``True`` (default), descend into subdirectories.
    :param verbose:        If ``True``, log each file path as it is processed.
    :param ollama_url:     Ollama base URL for embedding and vision calls.
    :param caption_images: If ``True``, call the vision model on image-only
                           slides and replace their stub title with a real caption.
    """

    def __init__(
        self,
        index_dir:      Path  = DEFAULT_INDEX_DIR,
        collection:     str   = COLLECTION_NAME,
        batch_size:     int   = 32,
        force:          bool  = False,
        recursive:      bool  = True,
        verbose:        bool  = False,
        ollama_url:     str   = DEFAULT_OLLAMA_URL,
        caption_images: bool  = False,
    ) -> None:
        self._index_dir     = index_dir.expanduser().resolve()
        self._collection    = collection
        self._batch_size    = batch_size
        self._force         = force
        self._recursive     = recursive
        self._verbose       = verbose
        self._embed_url     = f"{ollama_url.rstrip('/')}/api/embed"
        self._caption_images = caption_images
        self._captioner     = SlideVisionCaptioner(ollama_url) if caption_images else None

        self._index_dir.mkdir(parents=True, exist_ok=True)

        self._client   = QdrantClient(path=str(self._index_dir / "qdrant"))
        self._manifest = MtimeManifest(self._index_dir / "manifest.json")

        self._ensure_collection()

    # ------------------------------------------------------------------
    def _ensure_collection(self) -> None:
        """Create the Qdrant collection if it does not already exist."""
        existing = {c.name for c in self._client.get_collections().collections}
        if self._collection not in existing:
            self._client.create_collection(
                collection_name=self._collection,
                vectors_config=VectorParams(size=EMBEDDING_DIM, distance=Distance.COSINE),
            )
            log.info(f"Created Qdrant collection '{self._collection}'.")

    # ------------------------------------------------------------------
    def index_inputs(self, inputs: list[Path]) -> None:
        """Index all presentations reachable from *inputs*.

        Iterates over discovered files, skips unchanged ones (unless
        ``--force``), extracts slide chunks, embeds them in batches, and
        upserts into Qdrant.  After indexing, prunes any files that are
        tracked in the manifest but no longer exist on disk.

        :param inputs: List of file or directory paths supplied on the CLI.
        """
        all_files   = list(self._discover(inputs))
        current_set = {str(p.resolve()) for p in all_files}

        total_files   = 0
        total_slides  = 0
        total_skipped = 0

        pending_points: list[PointStruct] = []

        for pres_path in all_files:
            abs_str = str(pres_path.resolve())

            if not self._force and not self._manifest.is_stale(pres_path):
                total_skipped += 1
                continue

            if self._verbose:
                log.info(f"Indexing: {pres_path}")

            try:
                chunks = self._extract(pres_path)
            except Exception as exc:
                log.warn(f"Skipping '{pres_path}': {exc}")
                continue

            if not chunks:
                log.warn(f"No text extracted from '{pres_path}' — skipping.")
                continue

            # Delete any previously indexed points for this file before
            # upserting the fresh set, so stale slide points are replaced.
            self._delete_file_points(abs_str)

            for chunk in chunks:
                vec = self._embed(chunk["text"])
                point_id = self._stable_id(abs_str, chunk["slide_index"])
                pending_points.append(
                    PointStruct(id=point_id, vector=vec, payload=chunk)
                )

            if len(pending_points) >= self._batch_size:
                self._flush(pending_points)
                pending_points = []

            self._manifest.mark(pres_path)
            total_files  += 1
            total_slides += len(chunks)

        # Flush remainder
        if pending_points:
            self._flush(pending_points)

        # Prune deleted files
        pruned = self._prune(current_set)

        self._manifest.save()
        log.info(
            f"Done. Indexed {total_files} file(s), {total_slides} slide(s). "
            f"Skipped {total_skipped} unchanged. Pruned {pruned} deleted."
        )

    # ------------------------------------------------------------------
    def _discover(self, inputs: list[Path]) -> Generator[Path, None, None]:
        """Yield all supported presentation files reachable from *inputs*.

        Handles a mix of individual files and directories.  When *recursive*
        is ``True`` (the default), descends into all subdirectories.

        :param inputs: List of file or directory paths.
        :return: Generator of ``Path`` objects for each presentation file found.
        """
        for inp in inputs:
            if inp.is_file():
                if inp.suffix.lower() in SUPPORTED_SUFFIXES:
                    yield inp
                else:
                    log.warn(f"'{inp}' is not a supported format — skipping.")
            elif inp.is_dir():
                if self._recursive:
                    for suffix in SUPPORTED_SUFFIXES:
                        yield from sorted(inp.rglob(f"*{suffix}"))
                else:
                    for suffix in SUPPORTED_SUFFIXES:
                        yield from sorted(inp.glob(f"*{suffix}"))
            else:
                log.warn(f"'{inp}' does not exist — skipping.")

    # ------------------------------------------------------------------
    def _extract(self, pres_path: Path) -> list[dict]:
        """Dispatch extraction to the appropriate format handler.

        For Keynote files, passes the vision captioner into ``KeynoteExtractor``
        so captioning happens inside the temp dir context while images are live.

        :param pres_path: Path to a ``.key`` or ``.pptx`` file.
        :return: List of slide chunk dicts.
        :raises ValueError: If the file suffix is unrecognised.
        """
        suffix = pres_path.suffix.lower()
        if suffix == ".key":
            captioner = self._captioner if self._caption_images else None
            return KeynoteExtractor(
                pres_path,
                captioner=captioner,
                verbose=self._verbose,
            ).extract()
        elif suffix == ".pptx":
            return PptxExtractor(pres_path).extract()
        else:
            raise ValueError(f"Unsupported format: '{suffix}'")

    # ------------------------------------------------------------------
    def _embed(self, text: str) -> list[float]:
        """Embed *text* via the Ollama nomic-embed-text REST API.

        Uses the ``passage:`` prefix expected by nomic-embed-text's
        asymmetric retrieval mode.

        :param text: Text to embed.
        :return: Embedding vector as a list of floats.
        :raises RuntimeError: If the Ollama API call fails.
        """
        prefixed = f"passage: {text}"
        resp = requests.post(
            self._embed_url,
            json={"model": OLLAMA_EMBED_MODEL, "input": prefixed},
            timeout=30,
        )
        if resp.status_code != 200:
            raise RuntimeError(
                f"Ollama embed API returned {resp.status_code}: {resp.text}"
            )
        data = resp.json()
        # Ollama returns {"embeddings": [[...float...]]}
        return data["embeddings"][0]

    # ------------------------------------------------------------------
    def _flush(self, points: list[PointStruct]) -> None:
        """Upsert a batch of points into Qdrant.

        :param points: List of ``PointStruct`` objects to upsert.
        """
        self._client.upsert(collection_name=self._collection, points=points)
        log.info(f"Upserted {len(points)} point(s).")

    # ------------------------------------------------------------------
    def _delete_file_points(self, abs_str: str) -> None:
        """Delete all Qdrant points previously indexed from *abs_str*.

        Uses a payload filter so only points from that file are removed.

        :param abs_str: Absolute path string used as the ``file_path`` payload key.
        """
        self._client.delete(
            collection_name=self._collection,
            points_selector=Filter(
                must=[FieldCondition(key="file_path", match=MatchValue(value=abs_str))]
            ),
        )

    # ------------------------------------------------------------------
    @staticmethod
    def _stable_id(abs_str: str, slide_index: int) -> int:
        """Generate a stable integer point ID from file path and slide index.

        Uses a hash so IDs are deterministic across runs, enabling upsert
        semantics (re-indexing the same slide overwrites the old point).

        :param abs_str:     Absolute path string of the presentation file.
        :param slide_index: 0-based slide position.
        :return: Non-negative integer suitable for a Qdrant point ID.
        """
        raw = f"{abs_str}::{slide_index}"
        return int(hashlib.md5(raw.encode()).hexdigest(), 16) % (2**53)

    # ------------------------------------------------------------------
    def _prune(self, current_set: set[str]) -> int:
        """Remove manifest entries and Qdrant points for files no longer on disk.

        :param current_set: Set of absolute path strings found in the current run.
        :return: Number of files pruned.
        """
        pruned = 0
        for known in list(self._manifest.known_paths()):
            if known not in current_set and not Path(known).exists():
                log.info(f"Pruning deleted file: {known}")
                self._delete_file_points(known)
                self._manifest.remove(Path(known))
                pruned += 1
        return pruned


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def _build_parser() -> argparse.ArgumentParser:
    """Build and return the argument parser.

    :return: Configured :class:`argparse.ArgumentParser` instance.
    """
    p = argparse.ArgumentParser(
        prog="presentation_indexer",
        description=(
            "Index Keynote (.key) and PowerPoint (.pptx) presentations "
            "into a local Qdrant vector store for semantic search."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    p.add_argument(
        "inputs",
        nargs="+",
        metavar="PATH",
        help="One or more .key/.pptx files or directories to index.",
    )
    p.add_argument(
        "--index-dir",
        type=Path,
        default=DEFAULT_INDEX_DIR,
        metavar="DIR",
        help=f"Qdrant storage + manifest directory. Default: {DEFAULT_INDEX_DIR}",
    )
    p.add_argument(
        "--collection",
        default=COLLECTION_NAME,
        metavar="NAME",
        help=f"Qdrant collection name. Default: {COLLECTION_NAME}",
    )
    p.add_argument(
        "--batch-size",
        type=int,
        default=32,
        metavar="N",
        help="Embedding batch size. Default: 32",
    )
    p.add_argument(
        "--force",
        action="store_true",
        help="Ignore mtime manifest and re-index all files.",
    )
    p.add_argument(
        "--no-recursive",
        action="store_true",
        help="Limit directory search to top level only (no subdirectory descent).",
    )
    p.add_argument(
        "--ollama-url",
        default=DEFAULT_OLLAMA_URL,
        metavar="URL",
        help=f"Ollama base URL. Default: {DEFAULT_OLLAMA_URL}",
    )
    p.add_argument(
        "--caption-images",
        action="store_true",
        help=(
            "Call the vision model (llama3.2-vision via Ollama) on image-only "
            "slides to generate searchable captions. Slow but improves recall "
            "for slides with no text. Requires Ollama to be running."
        ),
    )
    p.add_argument(
        "--verbose",
        action="store_true",
        help="Log each file path as it is processed.",
    )
    return p


def main() -> None:
    """Entry point for the presentation indexer CLI."""
    parser = _build_parser()
    args   = parser.parse_args()

    # Validate inputs; warn on missing paths but continue with the rest
    inputs: list[Path] = []
    for item in args.inputs:
        p = Path(item).expanduser()
        if not p.exists():
            log.warn(f"'{item}' does not exist — skipping.")
            continue
        inputs.append(p)

    if not inputs:
        parser.error("No valid inputs found.")

    indexer = PresentationIndexer(
        index_dir=args.index_dir.expanduser().resolve(),
        collection=args.collection,
        batch_size=args.batch_size,
        force=args.force,
        recursive=not args.no_recursive,
        verbose=args.verbose,
        ollama_url=args.ollama_url,
        caption_images=args.caption_images,
    )
    indexer.index_inputs(inputs)


if __name__ == "__main__":
    main()
